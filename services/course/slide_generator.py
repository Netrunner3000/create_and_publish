from __future__ import annotations
import os
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from .models import Slide, CourseOutline

# Slide canvas dimensions (16:9)
W, H = 1920, 1080

# Colour palette — professional dark theme
BG_COLOR = (15, 23, 42)          # slate-900
ACCENT_COLOR = (99, 102, 241)    # indigo-500
TITLE_COLOR = (248, 250, 252)    # slate-50
BULLET_COLOR = (203, 213, 225)   # slate-300
FOOTER_COLOR = (100, 116, 139)   # slate-500
SLIDE_NUM_COLOR = (71, 85, 105)  # slate-600

FONT_PATHS = [
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
]

FONT_PATHS_BOLD = [
    "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
]


def _load_font(paths: list[str], size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for p in paths:
        if os.path.exists(p):
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def _wrap_text(text: str, font: ImageFont.FreeTypeFont, max_width: int, draw: ImageDraw.Draw) -> list[str]:
    words = text.split()
    lines, current = [], ""
    for word in words:
        test = f"{current} {word}".strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def render_slide_image(
    slide: Slide,
    slide_number: int,
    total_slides: int,
    course_title: str,
    output_path: str,
) -> str:
    img = Image.new("RGB", (W, H), BG_COLOR)
    draw = ImageDraw.Draw(img)

    font_title = _load_font(FONT_PATHS_BOLD, 72)
    font_bullet = _load_font(FONT_PATHS, 42)
    font_footer = _load_font(FONT_PATHS, 28)
    font_slide_num = _load_font(FONT_PATHS, 26)

    # Accent bar top
    draw.rectangle([(0, 0), (W, 8)], fill=ACCENT_COLOR)

    # Course title footer
    draw.rectangle([(0, H - 70), (W, H)], fill=(10, 15, 30))
    draw.text((60, H - 50), course_title, font=font_footer, fill=FOOTER_COLOR, anchor="lm")

    # Slide number
    slide_num_text = f"{slide_number} / {total_slides}"
    draw.text((W - 60, H - 50), slide_num_text, font=font_slide_num, fill=SLIDE_NUM_COLOR, anchor="rm")

    # Accent line under title area
    draw.rectangle([(60, 160), (W - 60, 164)], fill=ACCENT_COLOR)

    # Title
    title_lines = _wrap_text(slide.title, font_title, W - 120, draw)
    y = 60
    for line in title_lines[:2]:
        draw.text((60, y), line, font=font_title, fill=TITLE_COLOR)
        bbox = draw.textbbox((0, 0), line, font=font_title)
        y += (bbox[3] - bbox[1]) + 8

    # Bullets
    y = 220
    line_height = 70
    max_content_height = H - 160  # stay above footer
    bullet_font = _load_font(FONT_PATHS, 42)

    for bullet in slide.bullets:
        if y + line_height > max_content_height:
            break
        bullet_lines = _wrap_text(f"• {bullet}", bullet_font, W - 160, draw)
        for i, bl in enumerate(bullet_lines):
            indent = 0 if i == 0 else 30
            draw.text((80 + indent, y), bl, font=bullet_font, fill=BULLET_COLOR)
            y += line_height
        y += 10  # extra gap between bullets

    img.save(output_path, "PNG", quality=95)
    return output_path


def render_title_slide(course: CourseOutline, output_path: str) -> str:
    img = Image.new("RGB", (W, H), BG_COLOR)
    draw = ImageDraw.Draw(img)

    font_title = _load_font(FONT_PATHS_BOLD, 96)
    font_subtitle = _load_font(FONT_PATHS, 52)
    font_meta = _load_font(FONT_PATHS, 36)

    # Full-height accent bar left
    draw.rectangle([(0, 0), (12, H)], fill=ACCENT_COLOR)
    # Bottom accent bar
    draw.rectangle([(0, H - 8), (W, H)], fill=ACCENT_COLOR)

    # Vertical centre
    center_y = H // 2 - 80

    title_lines = _wrap_text(course.title, font_title, W - 200, draw)
    y = center_y - len(title_lines) * 110 // 2
    for line in title_lines:
        draw.text((W // 2, y), line, font=font_title, fill=TITLE_COLOR, anchor="mm")
        bbox = draw.textbbox((0, 0), line, font=font_title)
        y += (bbox[3] - bbox[1]) + 16

    draw.rectangle([(W // 2 - 200, y + 20), (W // 2 + 200, y + 24)], fill=ACCENT_COLOR)

    y += 50
    draw.text((W // 2, y), course.subtitle, font=font_subtitle, fill=BULLET_COLOR, anchor="mm")
    y += 80
    meta = f"{course.difficulty.value.capitalize()}  ·  {course.total_hours:.1f} hours  ·  {course.target_audience}"
    draw.text((W // 2, y), meta, font=font_meta, fill=FOOTER_COLOR, anchor="mm")

    img.save(output_path, "PNG", quality=95)
    return output_path


def generate_lesson_slides(
    slides: list[Slide],
    course_title: str,
    lesson_dir: str,
) -> list[str]:
    os.makedirs(lesson_dir, exist_ok=True)
    paths = []
    for i, slide in enumerate(slides):
        out = os.path.join(lesson_dir, f"slide_{i+1:03d}.png")
        render_slide_image(slide, i + 1, len(slides), course_title, out)
        paths.append(out)
    return paths


def generate_pptx(
    slides: list[Slide],
    course_title: str,
    lesson_title: str,
    output_path: str,
) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        blank_layout = prs.slide_layouts[6]

        def hex_to_rgb(h: int, g: int, b: int) -> RGBColor:
            return RGBColor(h, g, b)

        for slide_data in slides:
            slide = prs.slides.add_slide(blank_layout)
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = hex_to_rgb(15, 23, 42)

            # Title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(15), Inches(1.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(248, 250, 252)

            # Bullets
            bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(14.5), Inches(6))
            tf2 = bullet_box.text_frame
            tf2.word_wrap = True
            for i, bullet in enumerate(slide_data.bullets):
                p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(26)
                p.font.color.rgb = hex_to_rgb(203, 213, 225)

            # Speaker notes
            slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

        prs.save(output_path)
    except ImportError:
        open(output_path, "w").write(f"# {lesson_title}\n\nInstall python-pptx to generate PPTX files.\n")

    return output_path
